import streamlit as st
from pptx import Presentation
from pptx.enum.lang import MSO_LANGUAGE_ID

from docx import Document

from ibm_watson_machine_learning.foundation_models.utils.enums import ModelTypes
from ibm_watson_machine_learning.foundation_models import Model
from ibm_watson_machine_learning.metanames import GenTextParamsMetaNames as GenParams
from ibm_watson_machine_learning.foundation_models.extensions.langchain import WatsonxLLM

import tempfile
import pathlib
import json
import ast

import os
from dotenv import load_dotenv

load_dotenv()

api_key = st.secrets["API_KEY"]
project_id = st.secrets["PROJECT_ID"]

api_key = os.getenv("API_KEY", None)
project_id = os.getenv("PROJECT_ID", None)

creds = {
    "url"    : "https://us-south.ml.cloud.ibm.com",
    "apikey" : api_key
}

params = {
    GenParams.DECODING_METHOD:"greedy",
    GenParams.MAX_NEW_TOKENS:3000,
    GenParams.MIN_NEW_TOKENS:1,
    GenParams.TOP_K:50,
    GenParams.TOP_P:1,
    # GenParams.STOP_SEQUENCES:["\\n\\n"],
}

llm = Model("meta-llama/llama-2-13b-chat",creds, params,project_id)

def buildprompt(text,sourcelang,targetlang):
    prompt = f"""[INST]you are a translator, and need to process a json file.
    -for each field in json, translate the field value in {sourcelang} language to {targetlang} language:
    -dont translate the field name.
    -dont show note.
    -show the json only.
    [/INST]
    {sourcelang} json:{text}
    {targetlang} json:"""
    # print(prompt)
    return prompt

def translate_doc(docfile,sourcelang,targetlang):
    adoc = Document(docfile)

    for paragraph in adoc.paragraphs:
        prompt = buildprompt(paragraph.text,sourcelang,targetlang)
        for response in llm.generate_text([prompt]):
            paragraph.text = response

    for table in adoc.tables:
        for row in table.rows:
            for cell in row.cells:
                prompt = buildprompt(cell.text,sourcelang,targetlang)
                for response in llm.generate_text([prompt]):
                    cell.text = response

    output_file_path = docfile.replace('.docx', f'-{targetlang}.docx')
    adoc.save(output_file_path)
    return output_file_path

def translate_ppt(pptfile,sourcelang,targetlang):
    # input_file_path = "sample.pptx"

    # sourcelang = MSO_LANGUAGE_ID.ENGLISH_US
    # targetlang = MSO_LANGUAGE_ID.CHINESE_HONG_KONG_SAR

    presentation = Presentation(pptfile)

    slide_number = 1
    for slide in presentation.slides:
        st.write(f'Slide {slide_number} of {len(presentation.slides)}')
        slide_number += 1

        context = {}

        # translate comments
        if slide.has_notes_slide:
            text_frame = slide.notes_slide.notes_text_frame
            if len(text_frame.text) > 0:
                context[text_frame.text] = text_frame.text

        # translate other texts
        for shape in slide.shapes:
            if shape.has_table:
                for cell in shape.table.iter_cells():
                    context[cell.text] = cell.text

            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for index, paragraph_run in enumerate(paragraph.runs):
                        context[paragraph_run.text] = paragraph_run.text

        # print(context)

        prompt = buildprompt(str(context),sourcelang,targetlang)
        allresp = ""
        for response in llm.generate_text([prompt]):
            print(response)
            allresp += response
        print(allresp)
        
        contextout = ast.literal_eval(allresp)

        # translate comments
        if slide.has_notes_slide:
            text_frame = slide.notes_slide.notes_text_frame
            if len(text_frame.text) > 0:
                slide.notes_slide.notes_text_frame.text = contextout[text_frame.text]


        # translate other texts
        for shape in slide.shapes:
            if shape.has_table:
                for cell in shape.table.iter_cells():
                    cell.text = contextout[cell.text]

            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for index, paragraph_run in enumerate(paragraph.runs):
                        paragraph.runs[index].text = contextout[paragraph_run.text]
                        # paragraph.runs[index].font.language_id = targetlang

    output_file_path = pptfile.replace('.pptx', f'-{targetlang}.pptx')
    presentation.save(output_file_path)
    return output_file_path

temp_dir = tempfile.TemporaryDirectory()

st.header("document translator powered by watsonx")

sourcelang = st.selectbox(
    'What language you want to translate from?',
    ('English','中文', 'Korean', 'Japanese','Thai','Malay','Bahasa Indonesia'))

targetlang = st.selectbox(
    'What language you want to translate to?',
    ('中文', 'Korean', 'Japanese','Thai','Malay','Bahasa Indonesia','English'))


uploaded_file = st.file_uploader(label="upload a ppt",type=['ppt','pptx','doc','docx'])
if uploaded_file is not None:
    st.write("filename:", uploaded_file.name)
    with open(os.path.join(pathlib.Path(temp_dir.name),uploaded_file.name),"wb") as f:
        f.write(uploaded_file.getbuffer())

    if uploaded_file.name.lower().endswith(('.ppt', '.pptx')):
        with st.spinner(text="In progress...", cache=False):
            output_file_path = translate_ppt(os.path.join(pathlib.Path(temp_dir.name),uploaded_file.name),sourcelang,targetlang)
        with open(output_file_path,'rb') as f:
            st.download_button('Download translated version', f,f'translate-{targetlang}.pptx')
    elif uploaded_file.name.lower().endswith(('.doc', '.docx')):
        with st.spinner(text="In progress...", cache=False):
            output_file_path = translate_doc(os.path.join(pathlib.Path(temp_dir.name),uploaded_file.name),sourcelang,targetlang)
        with open(output_file_path,'rb') as f:
            st.download_button('Download translated version', f,f'translate-{targetlang}.docx')