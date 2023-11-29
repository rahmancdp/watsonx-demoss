import streamlit as st
from pptx import Presentation
from pptx.enum.lang import MSO_LANGUAGE_ID

from docx import Document

from genai.credentials import Credentials
from genai.schemas import GenerateParams
from genai.model import Model

import tempfile
import pathlib

import os
from dotenv import load_dotenv

load_dotenv()

api_key = st.secrets["GENAI_KEY"]
api_endpoint = st.secrets["GENAI_API"]

api_key = os.getenv("GENAI_KEY", None)
api_endpoint = os.getenv("GENAI_API", None)

creds = Credentials(api_key,api_endpoint)

params = GenerateParams(
    decoding_method="greedy",
    max_new_tokens=1000,
    min_new_tokens=1,
    stream=False,
    top_k=50,
    top_p=1,
    stop_sequences=["\\n\\n"]
)

llm = Model(model="meta-llama/llama-2-70b-chat",credentials=creds, params=params)

def buildprompt(text,sourcelang,targetlang):
    return f"""[INST]be a translator, be concise.
    return the translated content only.
    please help translate following {sourcelang} in backquoted to {targetlang}:
    <<SYS>>
    {sourcelang}:`{text}`
    <<SYS>>
    [/INST]
    {targetlang}:"""

def translate_doc(docfile,sourcelang,targetlang):
    adoc = Document(docfile)

    for paragraph in adoc.paragraphs:
        prompt = buildprompt(paragraph.text,sourcelang,targetlang)
        for response in llm.generate([prompt]):
            paragraph.text = response.generated_text

    for table in adoc.tables:
        for row in table.rows:
            for cell in row.cells:
                prompt = buildprompt(cell.text,sourcelang,targetlang)
                for response in llm.generate([prompt]):
                    cell.text = response.generated_text

    output_file_path = docfile.replace('.docx', f'-{targetlang}.docx')
    adoc.save(output_file_path)
    return output_file_path

def translate_ppt(pptfile,sourcelang,targetlang):
    # input_file_path = "sample.pptx"

    # sourcelang = MSO_LANGUAGE_ID.ENGLISH_US
    # targetlang = MSO_LANGUAGE_ID.CHINESE_HONG_KONG_SAR

    presentation = Presentation(pptfile)

    slide_number = 1
    for slide in presentation.slides:
        st.write(f'Slide {slide_number} of {len(presentation.slides)}')
        slide_number += 1

        # translate comments
        if slide.has_notes_slide:
            text_frame = slide.notes_slide.notes_text_frame
            if len(text_frame.text) > 0:
                prompt = buildprompt(text_frame.text,sourcelang,targetlang)
                for response in llm.generate([prompt]):
                    slide.notes_slide.notes_text_frame.text = response.generated_text


        # translate other texts
        for shape in slide.shapes:
            if shape.has_table:
                for cell in shape.table.iter_cells():
                    prompt = buildprompt(cell.text,sourcelang,targetlang)
                    for response in llm.generate_async([prompt]):
                        cell.text = response.generated_text

            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for index, paragraph_run in enumerate(paragraph.runs):
                        prompt = buildprompt(paragraph_run.text,sourcelang,targetlang)
                        for response in llm.generate([prompt]):
                            paragraph.runs[index].text = response.generated_text
                            # paragraph.runs[index].font.language_id = targetlang

    output_file_path = pptfile.replace('.pptx', f'-{targetlang}.pptx')
    presentation.save(output_file_path)
    return output_file_path

temp_dir = tempfile.TemporaryDirectory()

st.header("document translator powered by watsonx")

sourcelang = st.selectbox(
    'What language you want to translate from?',
    ('English','Chinese', 'Korean', 'Japanese','Thai','Malay','Bahasa Indonesia'))

targetlang = st.selectbox(
    'What language you want to translate to?',
    ('Chinese', 'Korean', 'Japanese','Thai','Malay','Bahasa Indonesia','English'))


uploaded_file = st.file_uploader(label="upload a ppt",type=['ppt','pptx','doc','docx'])
if uploaded_file is not None:
    st.write("filename:", uploaded_file.name)
    with open(os.path.join(pathlib.Path(temp_dir.name),uploaded_file.name),"wb") as f:
        f.write(uploaded_file.getbuffer())

    if uploaded_file.name.lower().endswith(('.ppt', '.pptx')):
        output_file_path = translate_ppt(os.path.join(pathlib.Path(temp_dir.name),uploaded_file.name),sourcelang,targetlang)
        with open(output_file_path,'rb') as f:
            st.download_button('Download translated version', f,f'translate-{targetlang}.pptx')
    elif uploaded_file.name.lower().endswith(('.doc', '.docx')):
        output_file_path = translate_doc(os.path.join(pathlib.Path(temp_dir.name),uploaded_file.name),sourcelang,targetlang)
        with open(output_file_path,'rb') as f:
            st.download_button('Download translated version', f,f'translate-{targetlang}.docx')